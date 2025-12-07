from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import openai
import base64
import httpx
import os
import uuid
import json

router = APIRouter()

class DiapositivasRequest(BaseModel):
    template: str
    tema: str
    num_diapositivas: int

def get_openai_client():
    USER = os.getenv("LLM_USER", "rag_user").strip("'")
    PASSWORD = os.getenv("LLM_PASSWORD", "").strip("'")
    BASE_URL = os.getenv("OPENAI_BASE_URL")

    credentials = f"{USER}:{PASSWORD}"
    encoded_credentials = base64.b64encode(credentials.encode()).decode()

    http_client = httpx.Client(verify=False, timeout=120.0)

    client = openai.OpenAI(
        base_url=BASE_URL,
        api_key=os.getenv("OPENAI_API_KEY", "api_key"),
        default_headers={
            "Authorization": f"Basic {encoded_credentials}"
        },
        http_client=http_client
    )

    return client

def generar_contenido_llm(tema, num_slides):
    """Genera contenido usando el LLM"""
    try:
        client = get_openai_client()

        prompt = f"""Crea contenido educativo para una presentacion de {num_slides} diapositivas sobre: {tema}

Para cada diapositiva proporciona:
- Un titulo descriptivo
- 3-4 puntos clave concisos

Responde en formato JSON:
{{
  "diapositivas": [
    {{
      "titulo": "Titulo de la diapositiva",
      "puntos": ["Punto 1", "Punto 2", "Punto 3"]
    }}
  ]
}}"""

        completion = client.chat.completions.create(
            model=os.getenv("OPENAI_MODEL", "openai/gpt-oss-20b"),
            messages=[
                {"role": "system", "content": "Eres un experto en educacion creando presentaciones educativas claras y concisas."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=2000
        )

        contenido_texto = completion.choices[0].message.content.strip()

        if contenido_texto.startswith("```json"):
            contenido_texto = contenido_texto[7:]
        if contenido_texto.startswith("```"):
            contenido_texto = contenido_texto[3:]
        if contenido_texto.endswith("```"):
            contenido_texto = contenido_texto[:-3]
        contenido_texto = contenido_texto.strip()

        return json.loads(contenido_texto)

    except Exception as e:
        print(f"Error generando contenido con LLM: {str(e)}")
        # Fallback a contenido generico
        return {
            "diapositivas": [
                {
                    "titulo": f"Seccion {i+1}: {tema}",
                    "puntos": [
                        f"Concepto importante {i+1}.1",
                        f"Concepto importante {i+1}.2",
                        f"Concepto importante {i+1}.3"
                    ]
                }
                for i in range(num_slides)
            ]
        }

def aplicar_template_clasico_azul(prs, slide, titulo_texto, puntos, es_portada=False):
    """Template: Clasico Azul - Estilo profesional"""
    if es_portada:
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = titulo_texto
        subtitle.text = "Presentacion generada por Cuali.ai"

        for shape in [title, subtitle]:
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.color.rgb = RGBColor(43, 82, 135)
    else:
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = titulo_texto
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(43, 82, 135)

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)

        textbox = shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        for punto in puntos:
            p = text_frame.add_paragraph()
            p.text = f"• {punto}"
            p.font.size = Pt(20)
            p.space_before = Pt(12)
            p.font.color.rgb = RGBColor(0, 0, 0)

def aplicar_template_moderno_naranja(prs, slide, titulo_texto, puntos, es_portada=False):
    """Template: Moderno Naranja - Estilo dinamico"""
    if es_portada:
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = titulo_texto
        subtitle.text = "Presentacion generada por Cuali.ai"

        for shape in [title, subtitle]:
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.color.rgb = RGBColor(255, 69, 0)
                paragraph.font.bold = True
    else:
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = titulo_texto
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 69, 0)

        left = Inches(1.5)
        top = Inches(2.2)
        width = Inches(7)
        height = Inches(4)

        textbox = shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame

        for i, punto in enumerate(puntos):
            p = text_frame.add_paragraph()
            p.text = f"{i+1}. {punto}"
            p.font.size = Pt(18)
            p.space_before = Pt(15)
            p.font.color.rgb = RGBColor(50, 50, 50)

def aplicar_template_elegante_marino(prs, slide, titulo_texto, puntos, es_portada=False):
    """Template: Elegante Azul Marino - Formal"""
    if es_portada:
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = titulo_texto
        subtitle.text = "Presentacion generada por Cuali.ai"

        for shape in [title, subtitle]:
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.color.rgb = RGBColor(0, 51, 102)
    else:
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = titulo_texto
        title_shape.text_frame.paragraphs[0].font.size = Pt(30)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

        left = Inches(1.2)
        top = Inches(2.5)
        width = Inches(7.5)
        height = Inches(4)

        textbox = shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame

        for punto in puntos:
            p = text_frame.add_paragraph()
            p.text = f"▸ {punto}"
            p.font.size = Pt(19)
            p.space_before = Pt(10)
            p.font.color.rgb = RGBColor(40, 40, 40)

def aplicar_template_minimalista_verde(prs, slide, titulo_texto, puntos, es_portada=False):
    """Template: Minimalista Verde - Limpio"""
    if es_portada:
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = titulo_texto
        subtitle.text = "Presentacion generada por Cuali.ai"

        for shape in [title, subtitle]:
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.color.rgb = RGBColor(0, 128, 128)
    else:
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = titulo_texto
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 128, 128)

        left = Inches(1)
        top = Inches(2.3)
        width = Inches(8)
        height = Inches(4.5)

        textbox = shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame

        for punto in puntos:
            p = text_frame.add_paragraph()
            p.text = punto
            p.font.size = Pt(18)
            p.space_before = Pt(14)
            p.font.color.rgb = RGBColor(60, 60, 60)
            p.level = 1

@router.post("/generar-diapositivas")
async def generar_diapositivas(request: DiapositivasRequest):
    try:
        print(f"Generando presentacion: {request.tema} con {request.num_diapositivas} diapositivas")

        # Generar contenido con LLM
        contenido = generar_contenido_llm(request.tema, request.num_diapositivas)

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        templates_map = {
            'default': aplicar_template_clasico_azul,
            'ucalgary': aplicar_template_moderno_naranja,
            'leiden': aplicar_template_elegante_marino,
            'cuerna': aplicar_template_minimalista_verde
        }

        aplicar_template = templates_map.get(request.template, aplicar_template_clasico_azul)

        # Portada
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        aplicar_template(prs, slide, request.tema, [], es_portada=True)

        # Diapositivas de contenido
        for diapositiva_data in contenido["diapositivas"]:
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            aplicar_template(prs, slide, diapositiva_data["titulo"], diapositiva_data["puntos"])

        # Guardar
        filename = f"presentacion_{uuid.uuid4().hex[:8]}.pptx"
        filepath = f"/tmp/{filename}"
        prs.save(filepath)

        print(f"Presentacion guardada: {filepath}")

        return FileResponse(
            filepath,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=f"presentacion_{request.tema.replace(' ', '_')}.pptx"
        )

    except Exception as e:
        print(f"Error generando PPTX: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))